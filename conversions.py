"""
File conversion utilities for GemDesk
Handles conversion of various file formats to text/CSV for Gemini processing
"""

import tempfile
import requests
import os
import base64
from bs4 import BeautifulSoup


def convert_xlsx_to_csv(xlsx_path):
    """Convert Excel file to CSV"""
    try:
        import pandas as pd
        excel_file = pd.ExcelFile(xlsx_path)
        csv_content = ""
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(xlsx_path, sheet_name=sheet_name)
            csv_content += f"# Sheet: {sheet_name}\n"
            csv_content += df.to_csv(index=False)
            csv_content += "\n\n"
        temp_csv = tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False)
        temp_csv.write(csv_content)
        temp_csv.close()
        return temp_csv.name
    except ImportError:
        raise Exception("pandas and openpyxl required. Install: pip install pandas openpyxl")


def convert_ods_to_csv(ods_path):
    """Convert ODS (OpenDocument Spreadsheet) to CSV"""
    try:
        import pandas as pd
        excel_file = pd.read_excel(ods_path, engine='odf', sheet_name=None)
        csv_content = ""
        for sheet_name, df in excel_file.items():
            csv_content += f"# Sheet: {sheet_name}\n"
            csv_content += df.to_csv(index=False)
            csv_content += "\n\n"
        temp_csv = tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False)
        temp_csv.write(csv_content)
        temp_csv.close()
        return temp_csv.name
    except ImportError:
        raise Exception("pandas and odfpy required. Install: pip install pandas odfpy")


def convert_docx_to_pdf(docx_path):
    """Convert DOCX to PDF using python-docx and reportlab"""
    try:
        from docx import Document
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
        from reportlab.lib.units import inch
        from PIL import Image
        import io
        
        doc = Document(docx_path)
        temp_pdf = tempfile.NamedTemporaryFile(mode='wb', suffix='.pdf', delete=False)
        pdf_doc = SimpleDocTemplate(temp_pdf.name, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        temp_image_files = []  # Track temp files for cleanup
        
        # Process paragraphs and images
        for para in doc.paragraphs:
            if para.text.strip():
                # Detect style
                style = styles['Normal']
                if para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    if level == '1':
                        style = styles['Heading1']
                    elif level == '2':
                        style = styles['Heading2']
                    else:
                        style = styles['Heading3']
                
                # Clean text
                text = para.text.replace('\x00', '')
                p = Paragraph(text, style)
                story.append(p)
                story.append(Spacer(1, 0.1*inch))
        
        # Extract and add inline images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    image_stream = io.BytesIO(image_data)
                    img = Image.open(image_stream)
                    
                    # Resize if too large (max width 6 inches)
                    max_width = 6 * inch
                    aspect = img.height / img.width
                    if img.width > max_width:
                        img_width = max_width
                        img_height = max_width * aspect
                    else:
                        img_width = img.width
                        img_height = img.height
                    
                    # Save to temp file for reportlab
                    temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                    img.save(temp_img.name, 'PNG')
                    temp_img.close()
                    temp_image_files.append(temp_img.name)
                    
                    # Add to PDF
                    rl_img = RLImage(temp_img.name, width=img_width, height=img_height)
                    story.append(rl_img)
                    story.append(Spacer(1, 0.2*inch))
                except Exception as e:
                    print(f"Image extraction error: {e}")
                    continue
        
        # Process tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text for cell in row.cells)
                if row_text.strip():
                    p = Paragraph(row_text, styles['Normal'])
                    story.append(p)
            story.append(Spacer(1, 0.2*inch))
        
        try:
            # Build PDF (this reads the temp images)
            pdf_doc.build(story)
            temp_pdf.close()
        finally:
            # Always clean up temp images after PDF build attempt
            for temp_img_path in temp_image_files:
                try:
                    if os.path.exists(temp_img_path):
                        os.unlink(temp_img_path)
                except:
                    pass
        
        return temp_pdf.name
    except ImportError:
        raise Exception("python-docx and reportlab required. Install: pip install python-docx reportlab")


def convert_pptx_to_pdf(pptx_path):
    """Convert PowerPoint to PDF preserving images and charts"""
    try:
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib.enums import TA_CENTER
        from PIL import Image
        import io
        
        prs = Presentation(pptx_path)
        temp_pdf = tempfile.NamedTemporaryFile(mode='wb', suffix='.pdf', delete=False)
        pdf_doc = SimpleDocTemplate(temp_pdf.name, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        temp_image_files = []  # Track temp files for cleanup AFTER PDF build
        
        # Create centered style for slide titles
        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading1'],
            alignment=TA_CENTER,
            spaceAfter=12
        )
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Slide number header
            story.append(Paragraph(f"<b>Slide {slide_num}</b>", title_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Extract text and images from shapes
            for shape in slide.shapes:
                # Handle text
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.replace('\x00', '')
                    story.append(Paragraph(text, styles['Normal']))
                    story.append(Spacer(1, 0.1*inch))
                
                # Handle all image types (pictures, shapes with fills, etc.)
                try:
                    # Method 1: Direct picture shapes (shape_type 13)
                    if hasattr(shape, 'image'):
                        image_bytes = shape.image.blob
                        image_stream = io.BytesIO(image_bytes)
                        img = Image.open(image_stream)
                        
                        # Resize if too large (max width 6 inches)
                        max_width = 6 * inch
                        aspect = img.height / img.width
                        if img.width > max_width:
                            img_width = max_width
                            img_height = max_width * aspect
                        else:
                            img_width = img.width
                            img_height = img.height
                        
                        # Save to temp file
                        temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                        img.save(temp_img.name, 'PNG')
                        temp_img.close()
                        temp_image_files.append(temp_img.name)
                        
                        # Add to PDF
                        rl_img = RLImage(temp_img.name, width=img_width, height=img_height)
                        story.append(rl_img)
                        story.append(Spacer(1, 0.2*inch))
                        continue
                except Exception as e:
                    pass  # Try next method
                
                # Method 2: Charts (rendered as images)
                try:
                    if hasattr(shape, 'chart'):
                        # Export chart by rendering the shape
                        # Charts in python-pptx are complex - try to get the image representation
                        chart_img_bytes = shape.chart.part.blob
                        image_stream = io.BytesIO(chart_img_bytes)
                        img = Image.open(image_stream)
                        
                        max_width = 6 * inch
                        aspect = img.height / img.width
                        img_width = max_width
                        img_height = max_width * aspect
                        
                        temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                        img.save(temp_img.name, 'PNG')
                        temp_img.close()
                        temp_image_files.append(temp_img.name)
                        
                        rl_img = RLImage(temp_img.name, width=img_width, height=img_height)
                        story.append(rl_img)
                        story.append(Spacer(1, 0.2*inch))
                        continue
                except Exception as e:
                    pass  # Try next method
                
                # Method 3: Group shapes (diagrams, SmartArt, etc.)
                try:
                    if hasattr(shape, 'shapes'):  # It's a group shape
                        # Recursively extract images from grouped shapes
                        for subshape in shape.shapes:
                            if hasattr(subshape, 'image'):
                                image_bytes = subshape.image.blob
                                image_stream = io.BytesIO(image_bytes)
                                img = Image.open(image_stream)
                                
                                max_width = 6 * inch
                                aspect = img.height / img.width
                                if img.width > max_width:
                                    img_width = max_width
                                    img_height = max_width * aspect
                                else:
                                    img_width = img.width
                                    img_height = img.height
                                
                                temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                                img.save(temp_img.name, 'PNG')
                                temp_img.close()
                                temp_image_files.append(temp_img.name)
                                
                                rl_img = RLImage(temp_img.name, width=img_width, height=img_height)
                                story.append(rl_img)
                                story.append(Spacer(1, 0.2*inch))
                except Exception as e:
                    pass  # Continue to next shape
            
            # Page break between slides
            if slide_num < len(prs.slides):
                story.append(PageBreak())
        
        try:
            # Build PDF (this reads the temp images)
            pdf_doc.build(story)
            temp_pdf.close()
        finally:
            # Always clean up temp images after PDF build attempt
            for temp_img_path in temp_image_files:
                try:
                    if os.path.exists(temp_img_path):
                        os.unlink(temp_img_path)
                except:
                    pass
        
        return temp_pdf.name
    except ImportError:
        raise Exception("python-pptx and reportlab required. Install: pip install python-pptx reportlab")


def convert_odp_to_text(odp_path):
    """Convert ODP (OpenDocument Presentation) to text"""
    try:
        from odf import text as odf_text
        from odf.opendocument import load
        
        doc = load(odp_path)
        text_content = ""
        
        for paragraph in doc.getElementsByType(odf_text.P):
            if paragraph.firstChild:
                text_content += str(paragraph.firstChild) + "\n"
        
        temp_txt = tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8')
        temp_txt.write(text_content)
        temp_txt.close()
        return temp_txt.name
    except ImportError:
        raise Exception("odfpy required. Install: pip install odfpy")


def convert_odt_to_text(odt_path):
    """Convert ODT (OpenDocument Text) to text"""
    try:
        from odf import text as odf_text
        from odf.opendocument import load
        
        doc = load(odt_path)
        text_content = ""
        
        for paragraph in doc.getElementsByType(odf_text.P):
            if paragraph.firstChild:
                text_content += str(paragraph.firstChild) + "\n"
        
        temp_txt = tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8')
        temp_txt.write(text_content)
        temp_txt.close()
        return temp_txt.name
    except ImportError:
        raise Exception("odfpy required. Install: pip install odfpy")


def download_file_from_url(url):
    """Download a file directly from URL (PDFs, images, videos, etc.)"""
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.get(url, headers=headers, timeout=30, stream=True)
        response.raise_for_status()
        
        # Get file extension from URL or content-type
        import mimetypes
        from urllib.parse import urlparse
        
        # Get content-type from response
        content_type = response.headers.get('content-type', '').split(';')[0].strip()
        print(f"DEBUG download_file_from_url: Content-Type from server: {content_type}")
        
        # Try to get extension from URL first
        path = urlparse(url).path
        filename = path.split('/')[-1] if path else ''
        
        # Only treat as extension if it's a known file type (not version numbers like .v1)
        valid_extensions = {'.pdf', '.png', '.jpg', '.jpeg', '.gif', '.webp', '.mp4', '.mov', 
                           '.avi', '.mp3', '.wav', '.doc', '.docx', '.xlsx', '.pptx', '.csv', '.txt', '.zip'}
        ext = ''
        if '.' in filename:
            potential_ext = filename[filename.rfind('.'):].lower()
            if potential_ext in valid_extensions:
                ext = potential_ext
        
        print(f"DEBUG download_file_from_url: Extension from URL: {ext}")
        
        # If no extension in URL, determine from content-type
        if not ext:
            # Map content-type to extension
            if content_type == 'application/pdf':
                ext = '.pdf'
            elif content_type.startswith('image/'):
                ext = mimetypes.guess_extension(content_type) or '.jpg'
            elif content_type.startswith('video/'):
                ext = mimetypes.guess_extension(content_type) or '.mp4'
            else:
                ext = mimetypes.guess_extension(content_type) or '.bin'
            print(f"DEBUG download_file_from_url: Determined extension: {ext}")
        
        # Download to temp file with correct extension
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        print(f"DEBUG download_file_from_url: Temp file will be: {temp_file.name}")
        for chunk in response.iter_content(chunk_size=8192):
            if chunk:
                temp_file.write(chunk)
        temp_file.close()
        
        return temp_file.name
    except Exception as e:
        raise Exception(f"Failed to download file from URL: {str(e)}")


def is_direct_file_url(url):
    """Check if URL points to a direct file (not HTML page)"""
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.head(url, headers=headers, timeout=10, allow_redirects=True)
        content_type = response.headers.get('content-type', '').lower()
        
        # Check if it's a file type we support
        file_types = [
            'application/pdf',
            'image/',
            'video/',
            'audio/',
            'application/vnd.ms-excel',
            'application/vnd.openxmlformats',
            'application/vnd.oasis.opendocument',
            'application/msword',
            'text/plain',
            'text/csv',
        ]
        
        for file_type in file_types:
            if file_type in content_type:
                return True
        
        # Also check URL extension
        url_lower = url.lower()
        file_extensions = ['.pdf', '.jpg', '.jpeg', '.png', '.gif', '.webp', '.mp4', '.mov', 
                          '.avi', '.mp3', '.wav', '.xlsx', '.docx', '.pptx', '.csv', '.txt']
        return any(url_lower.endswith(ext) for ext in file_extensions)
    except:
        return False


def scrape_url(url):
    """Scrape text content from a URL"""
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        
        for script in soup(["script", "style"]):
            script.decompose()
        
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        temp_txt = tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8')
        temp_txt.write(f"URL: {url}\n\n{text}")
        temp_txt.close()
        return temp_txt.name
    except Exception as e:
        raise Exception(f"Failed to scrape URL: {str(e)}")


def generate_thumbnail(file_path, mime_type, size=(100, 100)):
    """Generate thumbnail for image, PDF, or video files"""
    try:
        from PIL import Image
        import io
        
        if 'image' in mime_type:
            # Images - direct thumbnail
            img = Image.open(file_path)
            img.thumbnail(size, Image.Resampling.LANCZOS)
            
            # Convert to base64
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            return base64.b64encode(buffered.getvalue()).decode()
            
        elif 'pdf' in mime_type:
            # PDFs - extract first page
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(file_path)
                page = doc[0]
                pix = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5))  # 50% scale
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                img.thumbnail(size, Image.Resampling.LANCZOS)
                doc.close()
                
                buffered = io.BytesIO()
                img.save(buffered, format="PNG")
                return base64.b64encode(buffered.getvalue()).decode()
            except ImportError:
                return None  # PyMuPDF not installed
                
        elif 'video' in mime_type:
            # Videos - extract frame 100 (skip black intro frames)
            try:
                import cv2
                cap = cv2.VideoCapture(file_path)
                
                # Try to get frame 100, fallback to frame 30 if video is shorter
                cap.set(cv2.CAP_PROP_POS_FRAMES, 100)
                ret, frame = cap.read()
                
                # If frame 100 fails, try frame 30
                if not ret:
                    cap.set(cv2.CAP_PROP_POS_FRAMES, 30)
                    ret, frame = cap.read()
                
                cap.release()
                
                if ret:
                    # Convert BGR to RGB
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    img = Image.fromarray(frame_rgb)
                    img.thumbnail(size, Image.Resampling.LANCZOS)
                    
                    buffered = io.BytesIO()
                    img.save(buffered, format="PNG")
                    return base64.b64encode(buffered.getvalue()).decode()
                return None
            except ImportError:
                return None  # opencv not installed
        
        return None  # Unsupported type
        
    except Exception as e:
        print(f"Thumbnail generation error: {e}")
        return None
